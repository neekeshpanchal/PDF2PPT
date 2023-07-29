import os
import requests
import pdfplumber
import urllib
from bs4 import BeautifulSoup
from transformers import pipeline
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import tkinter as tk
from tkinter import filedialog, messagebox, Label
from tkinter.ttk import Style, Button, Combobox
from tkinter import ttk
from PIL import Image, ImageTk



def scrape_image_from_word(word):

    # Denote search expression
    search_query = "https://www.google.com/search?q={word}&tbm=isch"

    # Send an HTTP Request
    response = requests.get(search_query)

    # Parse the HTML with BS4
    soup = BeautifulSoup(response.text, 'html.parser')

    # Find the first image
    image_element = soup.find('img')


    image_url = image_element['src']

    # Download
    download_image(image_url, word)


def download_image(url, word):
    # Create a directory to store the images if it doesn't exist
    if not os.path.exists("images"):
        os.makedirs("images")

    # Build the file path to save the image
    image_path = os.path.join("images", f"{word}.jpg")

    # Download the image and save it
    urllib.request.urlretrieve(url, image_path)




def summarize_text(text, num_sentences=3):
    #Load the summarizing model
    new_sum = pipeline("summarization", model='Alred/t5-small-finetuned-summarization-cnn')

    #Summarize the ingested text
    summarized = new_sum(text)

    return (summarized[0])['summary_text']


def convert_pdf_to_pptx(pdf_path, pptx_path, title_text):
    # Load the PDF using pdfplumber
    pdf_path = "C" + pdf_path
    pptx_path = pptx_path.strip()
    with pdfplumber.open(pdf_path.strip()) as pdf:
        # Create a PowerPoint presentation
        presentation = Presentation()

        # Set presentation background color to beige
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(7.5)
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(245, 245, 220)

        # Add the title slide
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title = title_slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].runs[0].font.size = Pt(45)
        title.text_frame.paragraphs[0].runs[0].font.name = "Times New Roman"
        title.text_frame.paragraphs[0].runs[0].font.bold = True

        # Define the section labels and their corresponding slide layouts
        section_labels = ["Abstract", "Introduction", "Methods", "Discussion", "Results"]
        slide_layouts = [1, 1, 1, 1, 1]

        for section_label, slide_layout in zip(section_labels, slide_layouts):
            # Add a slide for the section
            slide = presentation.slides.add_slide(presentation.slide_layouts[slide_layout])

            # Set the section label as the slide title
            title = slide.shapes.title
            title.text = section_label
            title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            title.text_frame.paragraphs[0].runs[0].font.name = "Times New Roman"
            title.text_frame.paragraphs[0].runs[0].font.bold = False

            text = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if section_label in page_text:
                    section_start = page_text.index(section_label) + len(section_label)
                    text = page_text[section_start:]
                    break

            # Summarize the section text
            summary = summarize_text(text)

            # Add the summarized text as the slide content if it's not empty
            if summary:
                # Get the existing text box on the slide (assuming there's only one text box per slide)
                slide_shapes = slide.shapes
                content = None
                for shape in slide_shapes:
                    if shape.has_text_frame:
                        content = shape.text_frame

                # Add the summarized text to the existing text box
                if content:
                    temp = ''
                    for char in summary: 
                        temp += char 
                        if char == '.':
                            temp += "\n"


                    content.text += "\n\n" + temp  # Add a new paragraph for the summarized text

                # Set the font size and name of the slide content
                    for paragraph in content.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(18)
                            run.font.name = "Times New Roman"

        # Save the PowerPoint presentation
        apply_theme(presentation, pptx_path)
        presentation.save(pptx_path)

        # Show a success message
        messagebox.showinfo("Success", "Conversion completed successfully!")


def select_pdf_file():
    # Open a file dialog to select the PDF file
    pdf_path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
    if pdf_path:
        pdf_label.configure(text="Selected PDF: " + pdf_path)


def select_pptx_file():
    # Open a file dialog to select the PowerPoint presentation file
    pptx_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if pptx_path:
        pptx_label.configure(text="Selected PowerPoint: " + pptx_path)


def start_conversion():
    # Get the selected file paths
    pdf_path = pdf_label.cget("text")[15:]
    pptx_path = pptx_label.cget("text")[20:]

    # Extract the research paper title
    title = os.path.splitext(os.path.basename(pdf_path))[0]

    # Convert the PDF to PPTX
    convert_pdf_to_pptx(pdf_path, pptx_path, title)


# Create the GUI window
window = tk.Tk()
window.title("PDF2PPT - Research Paper to PowerPoint Converter")
window.geometry("800x450")
window.configure(bg="#4960a6")

# Style the GUI elements
style = Style()
style.configure("TButton", font=("Times New Roman", 12))
style.configure("TLabel", font=("Times New Roman", 12))

# Add the title label
def load_image(filename, width, height):
    img = Image.open(filename)
    img = img.resize((width, height), Image.ANTIALIAS)
    return ImageTk.PhotoImage(img)

# Load the image for the title
title_image = load_image("PDF2PPTX-logos.jpeg", 150, 125) 

# Add the title label with the image
title_label = Label(window, image=title_image, background="#4960a6")
title_label.pack(pady=10)
# Add a label for the selected PDF
pdf_label = Label(window, text="Select your PDF", background="#4960a6", fg='#f6af85', font=('Calibri', 13, 'bold'))
pdf_label.pack()

# Add an "Upload PDF" button
pdf_button = Button(window, text="Upload PDF", command=select_pdf_file)
pdf_button.pack(pady=10)

# Add a label for the selected PowerPoint presentation
pptx_label = Label(window, text="Select your powerpoint", background="#4960a6", fg='#f6af85',font=('Calibri', 13, 'bold'))
pptx_label.configure(bg="#4960a6")
pptx_label.pack()

# Add an "Upload PowerPoint" button
pptx_button = Button(window, text="Upload PowerPoint", command=select_pptx_file)
pptx_button.pack(pady=10)

# Function to apply the selected theme
def apply_theme(presentation, pptx_path):
    selected_theme = theme_var.get()

    # Set the PowerPoint theme based on the user's selection
    if selected_theme == "Default":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(245, 245, 220)
    elif selected_theme == "Dark Mode":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(50, 50, 50)
    elif selected_theme == "Bright Colors":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(255, 102, 0)
    elif selected_theme == "Minimalistic":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(240, 240, 240)

    elif selected_theme == "Emerald Green":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(80, 240, 120)

    elif selected_theme == "Pastel Dreams":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(166, 233, 255)

    elif selected_theme == "Candy Crush":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(247, 133, 255)

    elif selected_theme == "Golden Sunset":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(255,215,0)

    elif selected_theme == "Plushy Coral":
        presentation.slide_master.background.fill.solid()
        presentation.slide_master.background.fill.fore_color.rgb = RGBColor(245, 149, 192)

    # Save the PowerPoint presentation with the applied theme
    presentation.save(pptx_path)


# Create a label for the theme dropdown
theme_label = Label(window, text="Select your theme", background="#4960a6", fg='#f6af85', font=('Calibri', 13, 'bold'))
theme_label.pack(pady=10)

# List of available themes
themes = ["Default", "Dark Mode", "Bright Colors", "Minimalistic", "Emerald Green", "Pastel Dreams", "Candy Crush", "Golden Sunset", "Plushy Coral"]
          
# Create a variable to hold the selected theme
theme_var = tk.StringVar()
theme_var.set(themes[0])  # Set the default theme

# Create the dropdown widget for selecting themes
theme_dropdown = ttk.Combobox(window, values=themes, textvariable=theme_var)
theme_dropdown.pack(pady=5)

# Add a "Start Conversion" button
conversion_button = Button(window, text="Start Conversion", command=start_conversion)
conversion_button.pack(pady=20)

# Start the GUI event loop
window.mainloop()