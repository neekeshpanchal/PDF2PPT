from pypdf import PdfReader 
from pptx import Presentation
from pathlib import Path
import PyPDF2
import tkinter as tk
from tkinter import filedialog



#Loading Summarization model
#summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

#PDF Preparation
reader = PyPDF2.PdfReader("example.pdf")

#Indexing Principles
introduction_text = []
method_text = []
results_text = []
conclusion_text = []

paperlength = len(reader.pages)
equalsplit = (paperlength/4)

for pagenumber in range(paperlength):
    if len(introduction_text) < equalsplit:
        introduction_text.append(reader.pages[pagenumber].extract_text())

    elif len(method_text) < equalsplit:
        method_text.append(reader.pages[pagenumber].extract_text())

    elif len(results_text) < equalsplit: 
        results_text.append(reader.pages[pagenumber].extract_text())

    elif len(conclusion_text) < equalsplit: 
        conclusion_text.append(reader.pages[pagenumber].extract_text())

#PPT Initialization
prs = Presentation('test.pptx')
prs.save('test.pptx')

SLD_LAYOUT_TITLE_AND_CONTENT = 1

slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
slide = prs.slides.add_slide(slide_layout)


import tkinter as tk
from tkinter import filedialog
import PyPDF2
import os

def upload_file():
    file_path = filedialog.askopenfilename()
    if file_path:
        # Process the file here
        print("Uploaded file:", file_path)
        display_pdf(file_path)
    else:
        print("No file selected.")

def display_pdf(file_path):
    # Clear any previously displayed PDF
    for widget in pdf_frame.winfo_children():
        widget.destroy()

    # Load the PDF
    with open(file_path, "rb") as file:
        pdf_reader = PyPDF2.PdfReader(file)
        total_pages = len(pdf_reader.pages)

        # Create a canvas to display the PDF
        canvas = tk.Canvas(pdf_frame, width=600, height=800)
        canvas.pack()

        # Render each page of the PDF
        for page_num in range(total_pages):
            page = pdf_reader.pages[page_num]
            img_data = page.to_image().convert("RGB").tobytes("raw", "RGB")
            image = tk.PhotoImage(data=img_data)
            canvas.create_image(0, 0, anchor=tk.NW, image=image)
            canvas.update()

        # Adjust the canvas scroll region
        canvas.configure(scrollregion=canvas.bbox("all"))

# Create the main window
window = tk.Tk()
window.title("PDF2PPT")

# Create a title label
title_label = tk.Label(window, text="PDF Viewer", font=("Helvetica", 18, "bold"))
title_label.pack(pady=20)

# Create a button
button = tk.Button(window, text="Upload File", command=upload_file, font=("Helvetica", 14))
button.pack(pady=10)

# Create a frame to hold the PDF viewer
pdf_frame = tk.Frame(window)
pdf_frame.pack(pady=20)

# Run the main event loop
window.mainloop()